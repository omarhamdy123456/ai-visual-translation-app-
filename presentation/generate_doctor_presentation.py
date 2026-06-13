"""Generate supervisor / doctor presentation for Translate Studio thesis project."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Doctor_Supervisor_Presentation.pptx"

ACCENT = RGBColor(0x1A, 0x56, 0x8E)  # academic blue
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, ACCENT)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(8.6), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.6), Inches(1.2))
    stf = sub.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    sp.alignment = PP_ALIGN.LEFT


def _add_section_slide(prs: Presentation, section: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = section
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(8.7), Inches(0.8))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.15), Inches(8.5), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)

    if note:
        nbox = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(8.5), Inches(0.6))
        np = nbox.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(12)
        np.font.color.rgb = MUTED
        np.font.italic = True


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "Translate Studio",
        "Multimodal OCR & Translation System\n"
        "[Your Name]  ·  [Supervisor Name]\n"
        "[Department / University]  ·  May 2026",
    )

    _add_content_slide(
        prs,
        "Outline",
        [
            "Problem & motivation",
            "Research objectives",
            "System architecture & methodology",
            "OCR engines, preprocessing, and translation",
            "Application modes (text, image, video, handwriting)",
            "Experimental evaluation (CER / WER)",
            "Results, challenges, and improvements",
            "Conclusion & future work",
        ],
    )

    _add_content_slide(
        prs,
        "Problem & Motivation",
        [
            "Travelers, migrants, and field workers often face text they cannot read (signs, menus, forms, subtitles).",
            "Manual translation is slow; generic phone apps lack control over OCR quality and layout preservation.",
            "Scene text varies in font, size, language, lighting, and camera angle — hurting recognition accuracy.",
            "Need: an integrated pipeline that extracts text reliably, translates it, and optionally redraws results on the image or video.",
        ],
    )

    _add_content_slide(
        prs,
        "Research Objectives",
        [
            "Design a modular multimodal translation pipeline (text, image, video, handwriting).",
            "Compare and combine OCR backends for multilingual scene text (Arabic + English focus).",
            "Apply image preprocessing to improve OCR on real-world photos (glare, noise, skew).",
            "Measure OCR quality with standard metrics (CER, WER) on fair benchmark datasets.",
            "Deliver a working prototype (Streamlit UI) demonstrating end-to-end translation workflows.",
        ],
    )

    _add_content_slide(
        prs,
        "System Overview — Translate Studio",
        [
            "Web application (Streamlit) branded SkyTranslate in the UI.",
            "Input: typed text, uploaded/camera images, video files, or handwriting (PDF/photo).",
            "Processing: OCR (when needed) → language detection → machine translation → optional visual overlay.",
            "Output: translated text, annotated images with bounding boxes, translated photos/videos, export (PDF/JSON).",
            "Translation default: Google Cloud Translate API (online, high quality).",
        ],
    )

    _add_section_slide(prs, "Architecture & Methodology")

    _add_content_slide(
        prs,
        "High-Level Architecture",
        [
            "UI layer: Streamlit (`app.py`) — mode selection, language settings, results tabs.",
            "Orchestration: `ocr_translator.py` — preprocess, OCR, box merge, overlay rendering, video sampling.",
            "OCR workers: isolated subprocess for PaddleOCR (`paddle_ocr_worker.py`) — stable on Windows.",
            "Translation: Google Cloud Translate (batch + segment-level for OCR lines).",
            "Evaluation: `eval_experimental_results.py`, `benchmark_funsd_fullpage.py` — reproducible metrics.",
        ],
        note="Replace UI with REST API (FastAPI) is a natural production extension.",
    )

    _add_content_slide(
        prs,
        "Technology Stack",
        [
            "Language: Python 3.11+",
            "OCR: PaddleOCR 3.x, EasyOCR, TrOCR (Hugging Face hybrid for handwriting)",
            "CV / preprocessing: OpenCV, PIL, optional RIDNet denoising, CLAHE, perspective correction",
            "Translation: Google Cloud Translate; OpenAI for handwriting transcription & photo generation",
            "ML runtime: PyTorch, PaddlePaddle; GPU/CPU selectable in the app",
        ],
    )

    _add_content_slide(
        prs,
        "OCR Methodology",
        [
            "Detection: find text regions (quadrilateral bounding boxes).",
            "Recognition: read each region with a language-specific model.",
            "PaddleOCR: one pass per language (e.g. Arabic then English), results merged.",
            "Post-processing: confidence filtering, duplicate merge, reading-order sort, RTL handling for Arabic.",
            "Outputs: (text, bbox, language tag, confidence) per line or word.",
        ],
    )

    _add_content_slide(
        prs,
        "Image Preprocessing (Before OCR)",
        [
            "Resize to effective max dimension while preserving readability.",
            "Optional GPU pipeline: denoise (DnCNN / RIDNet), illumination correction, sharpen, CLAHE.",
            "Optional perspective dewarp for angled documents and screens.",
            "Optional OpenAI image enhance for glare-heavy night photos (API cost).",
            "Goal: improve contrast and stroke clarity without destroying fine text.",
        ],
    )

    _add_content_slide(
        prs,
        "Translation Methodology",
        [
            "Each OCR segment translated once; duplicates deduplicated across video frames.",
            "Source language inferred per segment (Arabic script vs Latin vs mixed signage).",
            "Google Cloud Translate: batched requests for speed; fallback to local NLLB if cloud unavailable.",
            "Arabic target: reshaping + BiDi for correct RTL display on overlays.",
            "Optional: generate translated photo (local inpaint + text draw, or OpenAI image edit).",
        ],
    )

    _add_section_slide(prs, "Application Modes")

    _add_content_slide(
        prs,
        "Four Input Modes",
        [
            "Text — direct paste/type → translate (no OCR).",
            "Image — upload or camera → OCR → translate → optional bounding-box view & translated photo.",
            "Video — sample frames → OCR per frame → translate subtitles → optional translated video export.",
            "Handwriting — OpenAI Vision transcription → translate → bilingual PDF/PNG export.",
        ],
    )

    _add_section_slide(prs, "Experimental Evaluation")

    _add_content_slide(
        prs,
        "Evaluation Design",
        [
            "Metrics: Character Error Rate (CER) and Word Error Rate (WER) via `jiwer`.",
            "Arabic: line-level crops from `merged_arabic_rec` validation set (fair labels, ≤10 chars per line).",
            "English: full-page FUNSD test split (50 pages) — realistic document OCR, not single-word crops.",
            "Avoid unfair benchmarks: word-only FUNSD crops and paragraph-length Qari labels excluded from training eval.",
            "Hardware: NVIDIA RTX 3050 Laptop GPU; reproducible scripts under `experimental_results/`.",
        ],
    )

    _add_content_slide(
        prs,
        "Results — Arabic Line OCR (PaddleOCR pretrained)",
        [
            "Dataset: merged_arabic_rec validation — 222 line samples.",
            "CER (mean): 0.19 — about 1 character wrong per 5 characters.",
            "WER (mean): 0.60 — word boundaries remain challenging.",
            "Exact match rate: 43%.",
            "Interpretation: usable for short Arabic line crops; fine-tuning on line data recommended for production.",
        ],
        note="Source: experimental_results/ocr_results.json",
    )

    _add_content_slide(
        prs,
        "Results — English Full-Page Forms (FUNSD)",
        [
            "Dataset: nielsr/funsd test split — 50 full pages.",
            "PaddleOCR full-page CER ≈ 0.43, WER ≈ 0.58 (project benchmark).",
            "EasyOCR on same split: CER ≈ 0.49, WER ≈ 0.77 (reference comparison).",
            "Full-page OCR is harder than line crops — layout density and small fields affect recall.",
            "Demonstrates need for preprocessing + appropriate engine selection per document type.",
        ],
        note="Source: experimental_results/funsd_fullpage_test.json",
    )

    _add_content_slide(
        prs,
        "Translation Quality (Text Mode)",
        [
            "Offline NLLB-200 benchmark (English → Arabic, small demo set): BLEU ≈ 41.8.",
            "Production app uses Google Cloud Translate for image/video/text modes (higher fluency).",
            "Segment-level source-language inference avoids wrong en→en or de→de skips on multilingual OCR.",
        ],
    )

    _add_section_slide(prs, "Challenges & Engineering Fixes")

    _add_content_slide(
        prs,
        "Key Challenges Addressed",
        [
            "Multilingual signage: run Arabic + English Paddle passes; merge without dropping script-specific boxes.",
            "Repeated menu labels: overlay dedupe now uses text + position (not text alone) so duplicate items stay visible.",
            "Single-digit fields (e.g. order count “3”): lone digits no longer filtered from annotated view.",
            "Windows stability: PaddleOCR in isolated worker avoids Paddle/Torch DLL conflicts.",
            "Night / glare photos: preprocessing + optional OpenAI enhance; detector threshold tuning via environment variables.",
        ],
    )

    _add_content_slide(
        prs,
        "Live Demo (Suggested Flow)",
        [
            "1. Image mode: upload a bilingual sign or menu → show Annotated + Text tabs.",
            "2. Select Arabic + English under OCR languages; enable preprocess if needed.",
            "3. Translate to Arabic (or English) → show translated photo overlay.",
            "4. Optional: video mode with sampled frames and subtitle translation.",
            "5. Show JSON export with structured boxes for API-style consumption.",
        ],
        note="Insert screenshots from your app Annotated tab before the defense.",
    )

    _add_section_slide(prs, "Conclusion")

    _add_content_slide(
        prs,
        "Conclusions",
        [
            "Built a complete multimodal translation prototype combining multiple OCR engines and cloud translation.",
            "Quantitative evaluation shows fair baseline accuracy (Arabic line CER 0.19; English form CER ~0.43).",
            "Modular design separates UI, OCR orchestration, and workers — suitable for API deployment.",
            "Real-world issues (multilingual layout, dedupe, small digits) identified and fixed through systematic testing.",
        ],
    )

    _add_content_slide(
        prs,
        "Future Work",
        [
            "Fine-tune Arabic Paddle recognition on line-level dataset (`merged_arabic_rec`).",
            "Upgrade Paddle worker to 3.x API consistently; REST API (FastAPI) with authentication.",
            "Field-level schema extraction for forms (name, date, amount) beyond plain OCR text.",
            "Document similarity endpoint for fraud/KYC use cases.",
            "Mobile deployment and on-device OCR for offline-first scenarios.",
        ],
    )

    _add_title_slide(
        prs,
        "Thank You",
        "Questions & Discussion\n\n[Your email]  ·  [GitHub / project path]",
    )

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
